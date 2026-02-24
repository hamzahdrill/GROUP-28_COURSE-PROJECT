"""
Generate polished presentation slides for Group 28's AD-HTC Biorefinery Dashboard.
Output: Group28_AD-HTC_Slides.pptx
Includes: Theory, system architecture, API/SQL stack, cycle analysis, results.
"""

import math
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ═══════════════════════════════════════════════════════════════
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
CYCLE_IMG  = os.path.join(SCRIPT_DIR, "cycle_schematic.png")

# Colours
BG    = RGBColor(0x11, 0x11, 0x11)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIM   = RGBColor(0x88, 0x88, 0x88)
BLUE  = RGBColor(0x4A, 0x80, 0xD0)
GREEN = RGBColor(0x40, 0xB0, 0x60)
RED   = RGBColor(0xE0, 0x60, 0x40)
GOLD  = RGBColor(0xF0, 0xD0, 0x60)
PANEL = RGBColor(0x1A, 0x1A, 0x1A)
HDBG  = RGBColor(0x1B, 0x3A, 0x5C)

# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════
def set_bg(slide):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = BG

def box(slide, l, t, w, h, text, sz=14, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return tf

def bullets(slide, l, t, w, h, items, sz=13, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = "Calibri"; p.space_after = Pt(6)
    return tf

def title_bar(slide, title, sub=""):
    box(slide, 0.5, 0.3, 9, 0.6, title, sz=22, color=WHITE, bold=True, font="Consolas")
    if sub:
        box(slide, 0.5, 0.85, 9, 0.4, sub, sz=11, color=DIM, font="Consolas")
    ln = slide.shapes.add_shape(1, Inches(0.5), Inches(1.15), Inches(9), Emu(12000))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x33, 0x33, 0x33)
    ln.line.fill.background()

def tbl(slide, headers, rows, top=1.8):
    cols = len(headers); n = 1 + len(rows)
    ts = slide.shapes.add_table(n, cols, Inches(0.5), Inches(top), Inches(9), Inches(0.35 * n))
    t = ts.table
    for i, h in enumerate(headers):
        c = t.cell(0, i); c.text = h
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = WHITE
            p.font.name = "Consolas"; p.alignment = PP_ALIGN.CENTER
        c.fill.solid(); c.fill.fore_color.rgb = HDBG
    for ri, rd in enumerate(rows):
        for ci, v in enumerate(rd):
            c = t.cell(ri+1, ci); c.text = str(v)
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(9); p.font.color.rgb = WHITE
                p.font.name = "Consolas"; p.alignment = PP_ALIGN.CENTER
            c.fill.solid()
            c.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A) if ri%2==0 else PANEL

# ═══════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# ── 1. TITLE ──────────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
box(s, 1, 0.5, 8, 0.5, "MEG 315 — APPLIED THERMODYNAMICS",
    sz=12, color=DIM, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
box(s, 1, 1.1, 8, 1.0, "AD-HTC Integrated\nBiorefinery Dashboard",
    sz=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
box(s, 1, 2.4, 8, 0.5, "Otto Cycle  •  Rankine Cycle  •  Heat Integration  •  FastAPI  •  SQLite",
    sz=13, color=BLUE, align=PP_ALIGN.CENTER, font="Consolas")
box(s, 1, 3.3, 8, 1.0,
    "Group 28\nDepartment of Mechanical Engineering\nUniversity of Lagos  •  February 2026",
    sz=11, color=DIM, align=PP_ALIGN.CENTER)

# ── 2. OVERVIEW ───────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "PROJECT OVERVIEW", "What we built")
bullets(s, 0.6, 1.4, 8.8, 3.5, [
    "Interactive dashboard simulating the AD-HTC biorefinery system",
    "Otto Cycle (biogas ICE) — T-s diagram with 4 state points",
    "Rankine Steam Cycle — h-s Mollier diagram with pump, boiler, turbine, condenser",
    "Heat integration analysis — CHP waste heat vs AD + HTC demand",
    "12 real biomass feedstocks with published data (stored in SQLite)",
    "FastAPI backend with 4 RESTful API endpoints",
    "SQLite database for feedstock data and calculation history",
    "Animated SVG process schematic with live calculated values",
])

# ── 3. THEORY: WHAT IS AD-HTC? ───────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "THEORY: WHAT IS AD-HTC?", "Background and motivation")
bullets(s, 0.6, 1.4, 4.2, 3.8, [
    "Anaerobic Digestion (AD):",
    "  • Bacteria decompose organic matter",
    "    without oxygen",
    "  • Inputs: biomass (manure, food waste, etc.)",
    "  • Outputs: biogas (CH₄ + CO₂) + digestate",
    "  • Operates at 35–40°C (mesophilic)",
    "",
    "Hydrothermal Carbonization (HTC):",
    "  • Wet biomass heated under pressure",
    "    (180–260°C, 10–40 bar)",
    "  • Converts digestate → hydrochar",
    "  • Hydrochar: solid biofuel ≈ lignite coal",
], sz=11, color=WHITE)
bullets(s, 5.2, 1.4, 4.6, 3.8, [
    "Why integrate them?",
    "",
    "1. AD biogas → fuel for gas engine (electricity)",
    "",
    "2. Engine exhaust → waste heat recovered",
    "",
    "3. Waste heat → steam for HTC reactor",
    "",
    "4. HTC converts digestate → valuable hydrochar",
    "",
    "5. Result: closed-loop biorefinery maximising",
    "   energy recovery from biomass",
], sz=11, color=BLUE)

# ── 4. CYCLE SCHEMATIC ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "AD-HTC FUEL-ENHANCED GAS POWER CYCLE", "System schematic (Energhx, 2025)")
if os.path.exists(CYCLE_IMG):
    s.shapes.add_picture(CYCLE_IMG, Inches(0.8), Inches(1.3), Inches(8.4), Inches(4.0))
else:
    box(s, 0.6, 1.5, 8.8, 3.5,
        "[Place cycle_schematic.png in project folder to auto-embed]\n\n"
        "Biomass → Homogenizer → AD (moisture-rich) + HTC (moisture-lean)\n"
        "AD → Biogas → Combustion → Compressor/Turbine → Electricity\n"
        "HTC Steam Cycle: Boiler ↔ Reactor (closed loop)\n"
        "Exhaust → Waste Heat Recovery → Steam Generation",
        sz=12, color=DIM, align=PP_ALIGN.CENTER)
box(s, 0.6, 4.9, 8.8, 0.4,
    "Source: Energhx Research Group, Faculty of Engineering, University of Lagos",
    sz=9, color=DIM, align=PP_ALIGN.CENTER, font="Consolas")

# ── 5. THEORY: THERMODYNAMIC CYCLES ──────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "THEORY: THERMODYNAMIC CYCLES", "Otto gas cycle and Rankine steam cycle")
bullets(s, 0.5, 1.4, 4.4, 3.8, [
    "GAS POWER CYCLE (Otto / Brayton)",
    "",
    "1→2: Isentropic compression",
    "     T₂ = T₁ × r^(γ−1)",
    "",
    "2→3: Constant-volume combustion",
    "     Q_in = Cp × (T₃ − T₂)",
    "",
    "3→4: Isentropic expansion (power)",
    "     T₄ = T₃ / r^(γ−1)",
    "",
    "4→1: Heat rejection (exhaust)",
    "     W_net = W_turb − W_comp",
], sz=10, color=WHITE)
bullets(s, 5.2, 1.4, 4.6, 3.8, [
    "HTC STEAM CYCLE (Rankine)",
    "",
    "A→B: Pump — raises pressure",
    "     w_pump = vf × ΔP × 100",
    "",
    "B→C: Boiler — waste heat → steam",
    "     Superheated at P = 60 bar",
    "",
    "C→D: HTC reactor — steam heats",
    "     digestate for carbonization",
    "",
    "D→A: Condenser — steam returns",
    "     to liquid, cycle repeats",
], sz=10, color=GOLD)

# ── 6. TECH STACK & ARCHITECTURE ──────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "TECHNOLOGY STACK & ARCHITECTURE", "Full-stack web application")
tbl(s,
    ["Layer", "Technology", "Purpose"],
    [
        ["Backend", "Python / FastAPI", "RESTful API for cycle calculations"],
        ["Database", "SQLite", "Feedstock storage + calculation history"],
        ["Validation", "Pydantic", "Input/output data validation"],
        ["Frontend", "HTML / CSS / JavaScript", "Interactive dashboard UI"],
        ["Charts", "Chart.js", "T-s, h-s, energy, mass, heat charts"],
        ["Styling", "Tailwind CSS (CDN)", "Responsive layout"],
        ["Schematic", "SVG (programmatic)", "Animated process flow diagram"],
        ["Fonts", "Google Fonts (DM Sans/Mono)", "Typography"],
    ],
    top=1.5
)

# ── 7. API ENDPOINTS ─────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "API ENDPOINTS", "RESTful API — main.py (FastAPI)")
tbl(s,
    ["Method", "Endpoint", "Description", "Returns"],
    [
        ["GET", "/", "Serve dashboard HTML", "HTMLResponse"],
        ["GET", "/api/feedstocks", "All 12 feedstocks from SQLite", "JSON array"],
        ["POST", "/api/analyze", "Run full cycle analysis", "JSON result + stores in DB"],
        ["GET", "/api/history", "Past calculations", "JSON array from DB"],
    ],
    top=1.5
)
box(s, 0.6, 3.6, 8.8, 1.2,
    "POST /api/analyze accepts:\n"
    "  feedstock_key, feed_rate, T1, compression_ratio, T3, mass_flow_air,\n"
    "  t_steam_C, p_cond, eta_turbine, eta_pump, htc_temp\n\n"
    "Returns: otto{}, rankine{}, mass_balance{}, heat{}, olr",
    sz=10, color=DIM, font="Consolas")

# ── 8. DATABASE SCHEMA ────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "DATABASE SCHEMA", "SQLite — database.py")
tbl(s,
    ["Table", "Column", "Type", "Description"],
    [
        ["feedstocks", "id", "INTEGER PK", "Auto-increment"],
        ["feedstocks", "key / name", "TEXT", "Identifier & display name"],
        ["feedstocks", "ts, vs", "REAL", "Total / volatile solids fraction"],
        ["feedstocks", "biogas_yield, ch4", "REAL", "m³/kg VS, methane fraction"],
        ["feedstocks", "htc_yield, htc_hhv", "REAL", "Hydrochar yield & HHV"],
        ["calculations", "id, timestamp", "INT, TEXT", "Record ID & ISO timestamp"],
        ["calculations", "feedstock, inputs", "TEXT, JSON", "Key & full input params"],
        ["calculations", "otto_results", "JSON", "All Otto cycle outputs"],
        ["calculations", "rankine_results", "JSON", "All Rankine cycle outputs"],
        ["calculations", "heat_balance", "REAL", "kW surplus/deficit"],
    ],
    top=1.5
)

# ── 9. FEEDSTOCK DATABASE ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "FEEDSTOCK DATABASE", "12 real biomass types stored in SQLite")
tbl(s,
    ["Feedstock", "TS%", "VS%", "Biogas (m³/kg VS)", "CH₄%", "HTC Yield"],
    [
        ["Cow Manure", "18", "80", "0.280", "60", "48%"],
        ["Food Waste", "25", "90", "0.550", "62", "45%"],
        ["Rice Husk", "90", "80", "0.250", "52", "58%"],
        ["Microalgae", "10", "85", "0.450", "65", "40%"],
        ["Wood Chips", "85", "92", "0.200", "55", "65%"],
        ["Corn Stover", "88", "85", "0.338", "55", "55%"],
        ["Sugarcane Bagasse", "50", "90", "0.310", "54", "52%"],
    ],
    top=1.5
)
box(s, 0.6, 4.6, 8.8, 0.4,
    "+ Sewage Sludge, Pig Manure, MSW, Cassava Peel, Palm EFB",
    sz=10, color=DIM, align=PP_ALIGN.CENTER, font="Consolas")

# ── 10. OTTO CYCLE RESULTS ───────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "OTTO CYCLE RESULTS", "T₁=298K, r=10, T₃=1200K, ṁ_air=50 kg/s (from API)")

cp_air=1.005; cp_gas=1.148; g=1.4
T1o=298; rv=10; T3o=1200
T2o=T1o*rv**(g-1); T4o=T3o/rv**(g-1)
Wc=cp_air*(T2o-T1o); Wt=cp_gas*(T3o-T4o)
Wn=Wt-Wc; Qin=cp_gas*(T3o-T2o); eta=Wn/Qin

tbl(s,
    ["Point", "Process", "T (K)", "s (kJ/kg·K)"],
    [
        ["1", "Intake", f"{T1o:.1f}", "1.000"],
        ["2", "Compression", f"{T2o:.1f}", "1.000"],
        ["3", "Combustion", f"{T3o:.1f}", f"{1.0+cp_gas*math.log(T3o/T2o):.3f}"],
        ["4", "Exhaust", f"{T4o:.1f}", f"{1.0+cp_gas*math.log(T3o/T2o):.3f}"],
    ],
    top=1.5
)
box(s, 0.6, 3.5, 8.8, 0.8,
    f"W_comp = {Wc:.1f}   W_turb = {Wt:.1f}   W_net = {Wn:.1f} kJ/kg   η = {eta*100:.1f}%\n"
    f"Electrical Power = {50*Wn:.0f} kW   |   Thermal Recovery = {50*cp_gas*(T4o-T1o)*0.45:.0f} kW",
    sz=14, color=BLUE, bold=True, font="Consolas", align=PP_ALIGN.CENTER)

# ── 11. RANKINE CYCLE RESULTS ─────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "RANKINE CYCLE RESULTS", "P=60 bar, T_steam=500°C, P_cond=0.06 bar (from API)")

hA=359.9; sA=1.145; sfg=7.533-1.145; hfg=2653.6-359.9
wp=0.001*(60-0.06)*100; hB=hA+wp/0.80
hC=3423.1; sC=6.8826
xs=min(1.0,(sC-1.145)/sfg); hDs=hA+xs*hfg
hD=hC-0.85*(hC-hDs)
wn_r=(hC-hD)-wp/0.80; eta_r=wn_r/(hC-hB)

tbl(s,
    ["Point", "Description", "h (kJ/kg)", "s (kJ/kg·K)"],
    [
        ["A", "Pump Inlet (sat. liquid)", f"{hA:.1f}", f"{sA:.4f}"],
        ["B", "Pump Outlet", f"{hB:.1f}", f"{sA:.4f}"],
        ["C", "Turbine Inlet (500°C)", f"{hC:.1f}", f"{sC:.4f}"],
        ["D", "Turbine Exit", f"{hD:.1f}", f"{1.145+((hD-hA)/hfg)*sfg:.4f}"],
    ],
    top=1.5
)
box(s, 0.6, 3.5, 8.8, 0.6,
    f"w_net = {wn_r:.1f} kJ/kg    η_rankine = {eta_r*100:.1f}%",
    sz=16, color=GOLD, bold=True, font="Consolas", align=PP_ALIGN.CENTER)

# ── 12. HEAT INTEGRATION ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "HEAT INTEGRATION", "CHP waste heat vs AD + HTC thermal demand")
bullets(s, 0.6, 1.4, 8.8, 1.5, [
    "CHP waste heat = ṁ_air × Cp × (T₄ − T₁) × 0.45  (45% recovery factor)",
    "AD demand = (feed kg/s) × Cp_water × (37°C − ambient)",
    "HTC demand = (digestate kg/s) × Cp × (T_HTC − 37°C)",
    "Heat balance = CHP output − AD demand − HTC demand",
], sz=13, color=WHITE)
box(s, 0.6, 3.2, 8.8, 0.6,
    "SURPLUS → System is self-sufficient (no external heating)\n"
    "DEFICIT → External heating required",
    sz=13, color=GOLD, font="Consolas", align=PP_ALIGN.CENTER)
box(s, 0.6, 4.2, 8.8, 0.5,
    "API test: Heat balance = +4555.89 kW (SURPLUS) for Cow Manure at 10 t/day",
    sz=11, color=GREEN, align=PP_ALIGN.CENTER, font="Consolas")

# ── 13. MASS BALANCE ─────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "MASS & ENERGY BALANCE", "Cow Manure at 10 t/day (from API)")
feedKg=10000; TS=0.18; VS=0.80
vs=feedKg*TS*VS; bg_v=vs*0.280; ch4=bg_v*0.60; bg_m=bg_v*1.2
dd=vs*0.40; hc=dd*0.48
tbl(s,
    ["Output", "Value", "Unit"],
    [
        ["Volatile solids", f"{vs:.0f}", "kg/day"],
        ["Biogas volume", f"{bg_v:.0f}", "m³/day"],
        ["Methane volume", f"{ch4:.0f}", "m³/day"],
        ["Biogas mass", f"{bg_m:.0f}", "kg/day"],
        ["Dry digestate", f"{dd:.0f}", "kg/day"],
        ["Hydrochar", f"{hc:.0f}", "kg/day"],
    ],
    top=1.5
)

# ── 14. DEMO SCREENSHOTS ─────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "LIVE DEMO", "Screenshots from running dashboard")
bullets(s, 0.6, 1.5, 8.8, 3.5, [
    "[Insert screenshot: Dashboard with KPI cards and sidebar]",
    "[Insert screenshot: T-s diagram (Otto) and h-s Mollier (Rankine)]",
    "[Insert screenshot: Energy & mass balance charts]",
    "[Insert screenshot: Heat integration chart with surplus/deficit badge]",
    "[Insert screenshot: Animated process schematic]",
    "[Insert screenshot: FastAPI Swagger docs at /docs]",
], sz=13, color=DIM)

# ── 15. HOW TO RUN ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "HOW TO RUN")
bullets(s, 0.6, 1.4, 8.8, 3.5, [
    "1.  Install dependencies:",
    "        pip install -r requirements.txt",
    "",
    "2.  Start the server:",
    "        python main.py",
    "",
    "3.  Open in browser:",
    "        http://localhost:8000",
    "",
    "4.  API docs (auto-generated):",
    "        http://localhost:8000/docs",
    "",
    "5.  Or open ad-htc-biomass-v2.html directly",
    "     (standalone mode, no API)",
], sz=12, color=WHITE)

# ── 16. CONCLUSION ────────────────────────────────────────────
s = prs.slides.add_slide(prs.slide_layouts[6]); set_bg(s)
title_bar(s, "CONCLUSION")
bullets(s, 0.6, 1.4, 8.8, 3.2, [
    "Full-stack AD-HTC biorefinery simulation",
    "FastAPI backend with 4 RESTful endpoints",
    "SQLite database — feedstock storage + calculation history",
    "Otto cycle (biogas ICE) with T-s diagram — η = 72.6%",
    "Rankine steam cycle with h-s Mollier diagram — η = 27.7%",
    "Heat integration: CHP waste heat vs AD + HTC demand",
    "12 biomass feedstocks with real published data",
    "Animated process schematic with live values",
    "Zero-installation option: standalone HTML also works",
], sz=14)
box(s, 0.6, 4.6, 8.8, 0.4,
    "Group 28  •  MEG 315  •  University of Lagos  •  February 2026",
    sz=11, color=DIM, align=PP_ALIGN.CENTER, font="Consolas")

# ═══════════════════════════════════════════════════════════════
out = os.path.join(SCRIPT_DIR, "Group28_AD-HTC_Slides_Final.pptx")
prs.save(out)
print(f"Slides saved: {out}")
